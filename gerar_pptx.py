from pptx import Presentation

slides = [
    {"title": "Bem-vindo ao MetalmaOS", "content": "Gestão moderna de ordens de serviço"},
    {"title": "Centralize informações essenciais", "content": "Interface intuitiva, responsiva e segura"},
    {"title": "Acompanhe tudo em tempo real", "content": "Status das ordens, clientes, colaboradores e metas"},
    {"title": "Gestão completa de OS", "content": "Cadastro, associação de clientes/produtos/colaboradores, controle de tempo"},
    {"title": "Gestão inteligente de clientes", "content": "Filtros, validação automática, exportação"},
    {"title": "Produtividade individual", "content": "Metas de horas, acompanhamento de desempenho"},
    {"title": "Controle total de produtos", "content": "Estoque, preços, percentual global"},
    {"title": "Relatórios completos", "content": "Produtividade, tempo, status, gráficos e exportação"},
    {"title": "Personalização e segurança", "content": "Permissões de acesso, modo escuro"},
    {"title": "Suporte dedicado", "content": "Documentação completa, plataforma em evolução"},
    {"title": "MetalmaOS – Gestão moderna, resultados reais", "content": "Fale com nosso time e leve sua operação para o próximo nível!"},
]

prs = Presentation()
for slide in slides:
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = slide["title"]
    s.placeholders[1].text = slide["content"]

prs.save("MetalmaOS-Video-Slides.pptx")
print("Arquivo MetalmaOS-Video-Slides.pptx gerado com sucesso!") 