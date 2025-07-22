from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Imagens de fundo (deve estar na raiz do projeto)
FUNDOS = [
    os.path.abspath("fundo1.jpg"),
    os.path.abspath("fundo2.jpg"),
    os.path.abspath("fundo3.jpg"),
]
# Logomarcas
LOGO_CLARA = os.path.abspath("src/assets/logo.png")
LOGO_ESCURA = os.path.abspath("src/assets/logo2.png")

# Cores da empresa (extraídas do logo)
AZUL = RGBColor(0, 120, 212)
BRANCO = RGBColor(255, 255, 255)
CINZA_CLARO = RGBColor(220, 220, 220)

slides = [
    {"title": "Bem-vindo ao MetalmaOS", "content": "Gestão moderna de ordens de serviço"},
    {"title": "Centralize informações essenciais", "content": "Interface intuitiva, responsiva e segura"},
    {"title": "Acompanhe tudo em tempo real", "content": "Status das ordens, clientes, colaboradores e metas"},
    {"title": "Gestão completa de OS", "content": "Cadastro, associação de clientes/produtos/colaboradores, controle de tempo"},
    {"title": "Gestão inteligente de clientes", "content": "Filtros, validação automática, exportação"},
    {"title": "Produtividade individual", "content": "Metas de horas, acompanhamento de desempenho"},
    {"title": "Controle total de produtos", "content": "Estoque, preços, percentual global"},
    {"title": "Relatórios completos", "content": "Produtividade, tempo, status, gráficos e exportação"},
    {"title": "Personalização e segurança", "content": "Permissões de acesso, modo escuro"},
    {"title": "Suporte dedicado", "content": "Documentação completa, plataforma em evolução"},
    {"title": "MetalmaOS – Gestão moderna, resultados reais", "content": "Fale com nosso time e leve sua operação para o próximo nível!"},
]

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]  # Slide em branco

for idx, slide in enumerate(slides):
    s = prs.slides.add_slide(blank_slide_layout)
    width = prs.slide_width
    height = prs.slide_height

    # Fundo: alterna entre as imagens
    fundo_path = FUNDOS[idx % len(FUNDOS)]
    if os.path.exists(fundo_path):
        s.shapes.add_picture(fundo_path, 0, 0, width, height)

    # Overlay escuro translúcido para contraste
    shape = s.shapes.add_shape(1, 0, 0, width, height)  # 1 = Rectangle
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 28, 40)
    fill.transparency = 0.3
    shape.line.fill.background()

    # Logomarca (clara ou escura conforme fundo)
    logo_path = LOGO_CLARA if idx % 2 == 0 else LOGO_ESCURA
    if os.path.exists(logo_path):
        s.shapes.add_picture(logo_path, Inches(3.2), Inches(0.2), height=Inches(1.1))

    # Título
    title_box = s.shapes.add_textbox(Inches(0.5), Inches(1.5), width - Inches(1), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.add_paragraph()
    p.text = slide["title"]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = AZUL
    p.alignment = PP_ALIGN.CENTER

    # Texto
    text_box = s.shapes.add_textbox(Inches(1), Inches(3), width - Inches(2), Inches(2))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p2 = text_frame.add_paragraph()
    p2.text = slide["content"]
    p2.font.size = Pt(28)
    p2.font.color.rgb = BRANCO
    p2.alignment = PP_ALIGN.CENTER

    # Rodapé
    footer_box = s.shapes.add_textbox(Inches(0.5), height - Inches(0.7), width - Inches(1), Inches(0.5))
    footer_frame = footer_box.text_frame
    p3 = footer_frame.add_paragraph()
    p3.text = "MetalmaOS – Sistema de Controle de Ordens de Serviço"
    p3.font.size = Pt(16)
    p3.font.color.rgb = CINZA_CLARO
    p3.alignment = PP_ALIGN.CENTER

prs.save("MetalmaOS-Video-Slides-PROFISSIONAL-v2.pptx")
print("Arquivo MetalmaOS-Video-Slides-PROFISSIONAL-v2.pptx gerado com sucesso!") 