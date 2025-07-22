from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Caminho do logo
LOGO_PATH = os.path.abspath("Logo_Metalma_PNG.png")

# Cores da empresa (exemplo: azul escuro e cinza)
BG_COLOR = RGBColor(36, 41, 61)  # Fundo escuro
TITLE_COLOR = RGBColor(0, 120, 212)  # Azul destaque
TEXT_COLOR = RGBColor(40, 40, 40)  # Cinza escuro

slides = [
    {"title": "Bem-vindo ao MetalmaOS", "content": "Gestão moderna de ordens de serviço"},
    {"title": "Centralize informações essenciais", "content": "Interface intuitiva, responsiva e segura"},
    {"title": "Acompanhe tudo em tempo real", "content": "Status das ordens, clientes, colaboradores e metas"},
    {"title": "Gestão completa de OS", "content": "Cadastro, associação de clientes/produtos/colaboradores, controle de tempo"},
    {"title": "Gestão inteligente de clientes", "content": "Filtros, validação automática, exportação"},
    {"title": "Produtividade individual", "content": "Metas de horas, acompanhamento de desempenho"},
    {"title": "Controle total de produtos", "content": "Estoque, preços, percentual global"},
    {"title": "Relatórios completos", "content": "Produtividade, tempo, status, gráficos e exportação"},
    {"title": "Personalização e segurança", "content": "Permissões de acesso, modo escuro"},
    {"title": "Suporte dedicado", "content": "Documentação completa, plataforma em evolução"},
    {"title": "MetalmaOS – Gestão moderna, resultados reais", "content": "Fale com nosso time e leve sua operação para o próximo nível!"},
]

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]  # Slide em branco

for slide in slides:
    s = prs.slides.add_slide(blank_slide_layout)
    # Fundo colorido (retângulo)
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    shape = s.shapes.add_shape(1, left, top, width, height)  # 1 = Rectangle
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    shape.line.fill.background()

    # Logo no topo
    if os.path.exists(LOGO_PATH):
        s.shapes.add_picture(LOGO_PATH, Inches(0.3), Inches(0.3), height=Inches(1.0))

    # Título
    title_box = s.shapes.add_textbox(Inches(0.5), Inches(1.5), width - Inches(1), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.add_paragraph()
    p.text = slide["title"]
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Texto
    text_box = s.shapes.add_textbox(Inches(1), Inches(3), width - Inches(2), Inches(2))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p2 = text_frame.add_paragraph()
    p2.text = slide["content"]
    p2.font.size = Pt(26)
    p2.font.color.rgb = TEXT_COLOR
    p2.alignment = PP_ALIGN.CENTER

    # Rodapé
    footer_box = s.shapes.add_textbox(Inches(0.5), height - Inches(0.7), width - Inches(1), Inches(0.5))
    footer_frame = footer_box.text_frame
    p3 = footer_frame.add_paragraph()
    p3.text = "MetalmaOS – Sistema de Controle de Ordens de Serviço"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TITLE_COLOR
    p3.alignment = PP_ALIGN.CENTER

prs.save("MetalmaOS-Video-Slides-PROFISSIONAL.pptx")
print("Arquivo MetalmaOS-Video-Slides-PROFISSIONAL.pptx gerado com sucesso!") 